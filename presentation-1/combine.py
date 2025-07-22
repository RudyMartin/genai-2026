from pptx import Presentation
import os

# Set the folder path where the slides are stored
folder_path = "/path/to/your/folder"  # <--- Replace this with your actual path

# Create a new presentation to combine into
combined = Presentation()

# Loop through all .pptx files in the folder
for filename in sorted(os.listdir(folder_path)):
    if filename.endswith(".pptx"):
        path = os.path.join(folder_path, filename)
        prs = Presentation(path)
        for slide in prs.slides:
            # Use blank layout for each slide being copied
            slide_layout = combined.slide_layouts[6]
            new_slide = combined.slides.add_slide(slide_layout)
            for shape in slide.shapes:
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

# Save the combined file
combined.save(os.path.join(folder_path, "Final_Combined_Presentation.pptx"))
print("Presentation merged successfully!")
